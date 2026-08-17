"""Build docs/mid_review_update.pptx -- the genuine mid-review checkpoint.

Deliberately scoped to what's actually done at this point in the timeline:
data pipeline complete, Experiment 1 (baseline) trained/evaluated. Experiment
2 (cross-manipulation), Experiment 3 (compression robustness), Grad-CAM
interpretability, the live demo, and final conclusions are held back for the
end-review deck (docs/end_review_presentation.pptx, built by
scripts/build_end_review_ppt.py) -- showing everything now would leave
nothing to demonstrate progress toward at the final review. The course
guideline itself only specifies a mid-review *date*, not required content,
so this follows standard two-stage review practice.

Same visual language as the end-review deck (navy title slide, white
content slides, cyan accent, Calibri) so the two read as one continuous
project across both reviews.

Usage: python -m scripts.build_mid_review_ppt
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W, SLIDE_H = Emu(12191695), Emu(6858000)
MARGIN = Emu(640080)
FONT = "Calibri"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
CYAN = RGBColor(0x22, 0xB8, 0xCF)
TEAL = RGBColor(0x0E, 0x74, 0x90)
TEXT = RGBColor(0x1A, 0x20, 0x2C)
MUTED = RGBColor(0x47, 0x55, 0x69)
MUTED_L = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)
GRID = RGBColor(0xE2, 0xE8, 0xF0)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TOTAL_SLIDES = 10
PAGE = [0]

TEAM = [
    ("Venugopalan Gangadharan", "CB.AI.U4AID25115"),
    ("Vipin Sudhakar", "CB.AI.U4AID25166"),
    ("Rithvik Arulprakash", "CB.AI.U4AID25148"),
    ("Harshith Kv", "CB.AI.U4AID25119"),
]


def new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return slide


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=None,
            line_spacing=None, space_after=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    if not runs or not isinstance(runs[0], list):
        runs = [runs]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = Pt(space_after)
        for text, size, bold, color in line:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def header(slide, eyebrow, title, dark=False):
    eyebrow_color = CYAN if dark else TEAL
    title_color = WHITE if dark else TEXT
    textbox(slide, MARGIN, Emu(384048), Emu(9144000), Emu(320040),
            [(eyebrow, 13, True, eyebrow_color)])
    textbox(slide, MARGIN, Emu(658368), Emu(10852800), Emu(731520),
            [(title, 28, True, title_color)])
    rect(slide, MARGIN, Emu(1353312), Emu(502920), Emu(41148), CYAN)


def footer(slide, title="Deepfake Detection — XceptionNet vs. CNN-ViT (23AID205)"):
    PAGE[0] += 1
    textbox(slide, MARGIN, Emu(6510528), Emu(7315200), Emu(274320),
            [(title, 10, False, MUTED_L)])
    textbox(slide, Emu(10972800), Emu(6510528), Emu(822960), Emu(274320),
            [(f"{PAGE[0]}/{TOTAL_SLIDES}", 10, False, MUTED_L)],
            align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, h, items, size=15, color=TEXT, bold=False,
            space_after=10, line_spacing=1.12):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = f"▸  {item}"
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def stat_panel(slide, x, y, w, h, label, stats):
    rect(slide, x, y, w, h, CARD_BG)
    pad = Emu(228600)
    textbox(slide, x + pad, y + Emu(180000), w - 2 * pad, Emu(320040),
            [(label, 11, True, TEAL)])
    row_h = Emu(int((h - Emu(500000)) / max(len(stats), 1)))
    cy = y + Emu(500000)
    for num, caption in stats:
        textbox(slide, x + pad, cy, w - 2 * pad, Emu(365760),
                [(num, 22, True, NAVY)])
        textbox(slide, x + pad, cy + Emu(340000), w - 2 * pad, Emu(320040),
                [(caption, 10.5, False, MUTED)])
        cy += row_h


def picture_slide(prs, eyebrow, title, caption, image_path, img_w, img_h,
                  stats_label=None, stats=None):
    slide = new_slide(prs)
    header(slide, eyebrow, title)
    textbox(slide, MARGIN, Emu(1691640), Emu(10515600), Emu(457200),
            [(caption, 13.5, False, MUTED)])
    slide.shapes.add_picture(image_path, MARGIN, Emu(2286000), width=img_w, height=img_h)
    if stats:
        gap = Emu(274320)
        stat_x = MARGIN + img_w + gap
        stat_w = SLIDE_W - MARGIN - stat_x
        stat_panel(slide, stat_x, Emu(2286000), stat_w, Emu(4114800),
                  stats_label or "RESULT", stats)
    footer(slide)
    return slide


def card(slide, x, y, w, h, header_text, header_color, items, item_color):
    rect(slide, x, y, w, h, CARD_BG)
    rect(slide, x, y, w, Emu(502920), header_color)
    textbox(slide, x + Emu(274320), y + Emu(118872), w - Emu(548640), Emu(320040),
            [(header_text, 14, True, WHITE)])
    bullets(slide, x + Emu(274320), y + Emu(731520), w - Emu(548640), h - Emu(800000),
            items, size=12.5, color=item_color, bold=True, space_after=8)


def table_slide(prs, eyebrow, title, subtitle, headers, rows, col_widths,
                highlight_rows=None, note=None):
    slide = new_slide(prs)
    header(slide, eyebrow, title)
    y = Emu(1691640)
    if subtitle:
        textbox(slide, MARGIN, y, Emu(10515600), Emu(400000),
                [(subtitle, 13.5, False, MUTED)])
        y += Emu(450000)
    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl_w = sum(col_widths)
    row_h = Emu(420000)
    tbl_h = row_h * n_rows
    gshape = slide.shapes.add_table(n_rows, n_cols, MARGIN, y, tbl_w, tbl_h)
    tbl = gshape.table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for c, htext in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = FONT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    highlight_rows = highlight_rows or set()
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xE7) if r in highlight_rows else (WHITE if r % 2 else CARD_BG)
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12.5)
            p.runs[0].font.bold = (c == 0)
            p.runs[0].font.color.rgb = TEXT
            p.runs[0].font.name = FONT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if note:
        textbox(slide, MARGIN, y + tbl_h + Emu(180000), Emu(10852800), Emu(600000),
                [(note, 12, False, MUTED)])
    footer(slide)
    return slide


def content_slide(prs, eyebrow, title, items, size=15.5, subtitle=None):
    slide = new_slide(prs)
    header(slide, eyebrow, title)
    y = Emu(1783080)
    if subtitle:
        textbox(slide, MARGIN, y, Emu(10852800), Emu(457200),
                [(subtitle, 14, False, MUTED)])
        y += Emu(550000)
    bullets(slide, MARGIN, y, Emu(10852800), Emu(6858000) - y - Emu(700000),
            items, size=size, color=TEXT, space_after=14)
    footer(slide)
    return slide


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    # 1. Title
    s = new_slide(prs, bg=NAVY)
    rect(s, 0, 0, Emu(146304), SLIDE_H, CYAN)
    textbox(s, MARGIN, Emu(2148840), Emu(8229600), Emu(365760),
            [("AIML COURSE PROJECT · 23AID205 · MID-REVIEW", 15, True, CYAN)])
    textbox(s, MARGIN, Emu(2468880), Emu(10332720), Emu(1920240),
            [[("Deepfake Detection:", 34, True, WHITE)],
             [("XceptionNet vs. CNN-ViT Hybrid on FaceForensics++", 34, True, WHITE)]])
    textbox(s, MARGIN, Emu(4343400), Emu(9601200), Emu(640000),
            [("Data pipeline complete, baseline results in. Cross-manipulation "
              "generalization, compression robustness, and final conclusions "
              "to follow at the final review.", 14, False, MUTED_L)])
    rect(s, MARGIN, Emu(5074920), Emu(457200), Emu(27432), CYAN)
    textbox(s, MARGIN, Emu(5989320), Emu(10852800), Emu(320040),
            [(" · ".join(name for name, _ in TEAM), 13, True, RGBColor(0xE5, 0xE9, 0xF0))])
    textbox(s, MARGIN, Emu(6263640), Emu(7315200), Emu(274320),
            [("AI & Data Science — Amrita Vishwa Vidyapeetham, Coimbatore", 11, False, MUTED_L)])
    PAGE[0] += 1

    # 2. Motivation, Problem & Objectives
    content_slide(prs, "INTRODUCTION", "Motivation, Problem & Objectives", [
        "Deepfakes have moved from research curiosity to a real misinformation/fraud "
        "vector -- but published detectors usually report one headline accuracy "
        "number, on the same manipulation type and compression level they trained on.",
        "Problem: that number is misleading in practice -- a detector at 99% in a "
        "paper can collapse the moment the test distribution shifts, and that "
        "failure is invisible until deployment.",
        "Objective 1: fairly compare XceptionNet (CNN baseline) vs. a CNN-ViT hybrid "
        "(ResNet + Transformer) on identical data/protocol -- does global "
        "self-attention actually improve generalization?",
        "Objective 2: measure same-distribution accuracy, full 4×4 cross-manipulation "
        "generalization, and compression robustness -- not just one number.",
        "Objective 3: explain *why* generalization succeeds or fails (Grad-CAM), and "
        "ship a working, testable live demo -- not only offline metrics.",
    ], size=15)

    # 3. Literature Review
    content_slide(prs, "LITERATURE REVIEW", "Directly Related Work", [
        "Rössler, A. et al. (2019). FaceForensics++: Learning to Detect Manipulated "
        "Facial Images. ICCV. — source of the dataset and the standard 4-method/"
        "3-compression benchmark protocol this project follows; also the paper "
        "XceptionNet is drawn from as the baseline detector.",
        "Chollet, F. (2017). Xception: Deep Learning with Depthwise Separable "
        "Convolutions. CVPR. — the CNN baseline architecture, still a standard "
        "reference detector for this exact task in the FF++ literature.",
        "Dosovitskiy, A. et al. (2021). An Image is Worth 16x16 Words: Transformers "
        "for Image Recognition at Scale. ICLR. — motivates the Transformer-encoder "
        "half of this project's CNN-ViT hybrid and its global self-attention "
        "mechanism.",
        "Selvaraju, R. R. et al. (2017). Grad-CAM: Visual Explanations from Deep "
        "Networks via Gradient-based Localization. ICCV. — the interpretability "
        "method used to explain why cross-manipulation generalization fails.",
        "Gap identified: most FF++-benchmark work reports same-distribution accuracy "
        "only; head-to-head CNN vs. attention-based generalization comparisons, "
        "paired with an interpretability layer, are comparatively under-studied — "
        "the specific gap this project addresses.",
    ], size=14)

    # 4. Methodology overview (pipeline)
    s = new_slide(prs)
    header(s, "METHODOLOGY", "End-to-End Pipeline")
    steps = [
        ("1", "Raw video", "FF++ c23, 5 folders: real + 4 manipulation methods"),
        ("2", "Face extraction", "MTCNN crop, 20 frames/video, margin 1.3x"),
        ("3", "Official splits", "video-level train/val/test — no identity leakage"),
        ("4", "Train", "AMP, class-balanced sampler, cosine LR, early stop"),
        ("5", "Evaluate", "frame + video-level Acc/F1/AUC on held-out test"),
        ("6", "Analyze", "cross-manip matrix, gap, Grad-CAM, live demo"),
    ]
    x = MARGIN
    card_w = Emu(1750000)
    gap = Emu(58000)
    for i, (n, t, d) in enumerate(steps):
        cx = x + i * (card_w + gap)
        rect(s, cx, Emu(2200000), card_w, Emu(3400000), CARD_BG)
        rect(s, cx, Emu(2200000), card_w, Emu(500000), NAVY)
        textbox(s, cx + Emu(120000), Emu(2320000), card_w - Emu(240000), Emu(320000),
                [(n, 15, True, CYAN)])
        textbox(s, cx + Emu(120000), Emu(2800000), card_w - Emu(240000), Emu(700000),
                [(t, 13.5, True, TEXT)])
        textbox(s, cx + Emu(120000), Emu(3450000), card_w - Emu(240000), Emu(1900000),
                [(d, 10.5, False, MUTED)])
    textbox(s, MARGIN, Emu(5850000), Emu(10852800), Emu(400000),
            [("Steps 1-3 (data) and step 4-5 for Experiment 1 are complete as of this "
              "review -- steps 4-6 continue for Experiments 2 & 3.", 11.5, False, MUTED)])
    footer(s)

    # 4. Dataset & Preprocessing
    content_slide(prs, "DATASET & PREPROCESSING", "Data: Source & Preparation — Complete", [
        "FaceForensics++ (FF++): 5,000 videos — 1,000 real (YouTube interviews) + "
        "4×1,000 manipulated (Deepfakes, Face2Face, FaceSwap, NeuralTextures), all "
        "sharing the same real-video identities. c23 compression, official FF++ "
        "folder layout (Kaggle mirror).",
        "Official video-level train/val/test splits used throughout — critical, "
        "since random frame-level splits would leak identity between train/test "
        "and inflate accuracy.",
        "Preprocessing: MTCNN face detection on GPU, 20 frames/video, 1.3x crop "
        "margin (captures blending-boundary artifacts), resized to 299×299 (native "
        "for XceptionNet, downsampled to 224 for the hybrid).",
        "Result: ~99,987 face crops across all 5,000 videos — the full dataset is "
        "preprocessed and ready; this stage is complete.",
        "Official c0/c40 access still pending — being handled with a documented "
        "workaround, to be shown at the final review.",
    ], size=14.5)

    # 5. EDA: outlier check
    picture_slide(prs, "EDA", "Data Quality: Outlier Check (IQR)",
                 "Per-video metadata (frame count, file size) checked via box-plot IQR "
                 "across all 5,000 videos before committing to full preprocessing.",
                 "analyze/outputs/outlier_boxplots.png", Emu(7406640), Emu(3648169),
                 stats_label="RESULT",
                 stats=[("5,000", "videos analyzed (real + 4 methods)"),
                        ("209", "flagged on frame count (4.2%)"),
                        ("326", "flagged on file size (6.5%)"),
                        ("502", "unique outliers overall (10.0%)")])

    # 6. Models comparison
    s = new_slide(prs)
    header(s, "ALGORITHM SELECTION", "Models: XceptionNet vs. CNN-ViT Hybrid")
    card_y, card_h = Emu(1750000), Emu(3000000)
    card_w = Emu(5195000)
    card(s, MARGIN, card_y, card_w, card_h, "XCEPTIONNET (BASELINE)", NAVY, [
        "Depthwise-separable CNN (Chollet 2017), ImageNet-pretrained",
        "2048-d pooled features → dropout → linear → logit",
        "Purely convolutional — local receptive fields only",
        "The standard reference detector for this exact task",
    ], TEAL)
    card(s, MARGIN + card_w + Emu(228600), card_y, card_w, card_h, "CNN-VIT HYBRID", TEAL, [
        "ResNet-50 (stages 1-3) → 14×14×1024 map → 512-d tokens + CLS",
        "6-layer, 8-head Transformer encoder over the tokens",
        "Tests: does global self-attention improve generalization?",
        "Same optimizer/schedule/data — architecture is the only variable",
    ], NAVY)
    strip_y = card_y + card_h + Emu(228600)
    rect(s, MARGIN, strip_y, Emu(10852800), Emu(18288), GRID)
    textbox(s, MARGIN, strip_y + Emu(180000), Emu(10852800), Emu(320040),
            [("SHARED TRAINING & EVALUATION PROTOCOL (fair comparison — only the architecture differs)", 11, True, TEAL)])
    textbox(s, MARGIN, strip_y + Emu(560000), Emu(10852800), Emu(900000),
            [("AMP training, AdamW, cosine LR, batch size 32 (GPU-benchmarked — "
              "XceptionNet silently throttles above this), early stopping (patience "
              "4 on val AUC). Metrics: Accuracy/F1/AUC at frame & video level (video "
              "= mean of frame probabilities).", 13, False, MUTED)])
    footer(s)

    # 7. Results: Experiment 1 baseline
    table_slide(prs, "RESULTS · EXPERIMENT 1 — COMPLETE", "Baseline Comparison (same-distribution, c23)",
               "Both models trained and tested on all 4 methods — the first of three "
               "planned experiments, complete as of this review.",
               ["Model", "Video Acc", "Video F1", "Video AUC", "Frame Acc", "Frame AUC"],
               [["XceptionNet", "97.29%", "0.983", "0.9958", "95.49%", "0.986"],
                ["CNN-ViT Hybrid", "94.71%", "0.967", "0.9836", "92.76%", "0.962"]],
               [Emu(2400000), Emu(1700000), Emu(1700000), Emu(1700000), Emu(1700000), Emu(1700000)],
               note="XceptionNet leads on every same-distribution metric here. Whether "
                    "this holds under cross-manipulation and compression shift is what "
                    "Experiments 2 & 3 will show at the final review.")

    # 8. Status & Plan to Completion
    s = new_slide(prs)
    header(s, "STATUS", "What's Done, What's Next")
    card(s, MARGIN, Emu(1783080), Emu(5195000), Emu(4114800), "DONE", GREEN, [
        "Full pipeline built: preprocessing, training, evaluation, analysis scripts",
        "Dataset acquired, quality-checked (IQR outlier analysis), and fully "
        "preprocessed — 99,987 face crops",
        "Both model architectures implemented and verified end-to-end",
        "GPU capacity benchmarked — found & fixed a silent throughput cliff above "
        "batch size 32",
        "Experiment 1 (baseline comparison): both models trained, evaluated, "
        "and analyzed",
    ], GREEN)
    card(s, MARGIN + Emu(5195000) + Emu(228600), Emu(1783080), Emu(5195000), Emu(4114800),
        "PLANNED FOR FINAL REVIEW", AMBER, [
        "Experiment 2: cross-manipulation generalization — full 4×4 matrix, "
        "both models",
        "Experiment 3: compression robustness under heavier video compression",
        "Interpretability analysis (Grad-CAM) explaining *why* generalization "
        "succeeds or fails",
        "Live demo (upload a video, get a real/fake prediction) and final "
        "conclusions",
    ], AMBER)
    footer(s)

    # 9. References + Contributors + repo
    s = new_slide(prs)
    header(s, "REFERENCES & CONTRIBUTORS", "Closing")
    bullets(s, MARGIN, Emu(1783080), Emu(10852800), Emu(1600000), [
        "Rössler, A. et al. (2019). FaceForensics++: Learning to Detect Manipulated Facial Images. ICCV.",
        "Chollet, F. (2017). Xception: Deep Learning with Depthwise Separable Convolutions. CVPR.",
        "Dosovitskiy, A. et al. (2021). An Image is Worth 16x16 Words: Transformers for Image Recognition at Scale. ICLR.",
        "Selvaraju, R. R. et al. (2017). Grad-CAM: Visual Explanations from Deep Networks via Gradient-based Localization. ICCV.",
    ], size=12.5, space_after=10)
    rect(s, MARGIN, Emu(3800000), Emu(10852800), Emu(18288), GRID)
    textbox(s, MARGIN, Emu(4000000), Emu(10852800), Emu(320040),
            [("CONTRIBUTORS", 11, True, TEAL)])
    bullets(s, MARGIN, Emu(4340000), Emu(10852800), Emu(1080000),
            [f"{name} — {roll}" for name, roll in TEAM],
            size=12, color=TEXT, bold=True, space_after=6)
    textbox(s, MARGIN, Emu(5500000), Emu(10852800), Emu(320040),
            [("AI & Data Science — Amrita Vishwa Vidyapeetham, Coimbatore", 11, False, MUTED)])
    textbox(s, MARGIN, Emu(5850000), Emu(10852800), Emu(320040),
            [("GITHUB REPOSITORY", 11, True, TEAL)])
    textbox(s, MARGIN, Emu(6100000), Emu(10000000), Emu(360000),
            [("github.com/DUNE-ODYSSEY/deepfake-detection", 12.5, True, CYAN)])
    footer(s)

    out = "docs/mid_review_update.pptx"
    prs.save(out)
    print(f"saved {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
